"""Generate group meeting presentation PPT for baseline algorithm correctness audit results."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os

# ── Color palette ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x24, 0x24, 0x3E)
ACCENT_BLUE = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_PURPLE = RGBColor(0xBB, 0x86, 0xFC)
ACCENT_GREEN = RGBColor(0x66, 0xBB, 0x6A)
TEXT_WHITE = RGBColor(0xF0, 0xF0, 0xF0)
TEXT_GRAY = RGBColor(0xB0, 0xB0, 0xC0)
TEXT_DIM = RGBColor(0x80, 0x80, 0x99)

# Method colors for charts
COLOR_OCC = RGBColor(0x42, 0xA5, 0xF5)     # Blue
COLOR_DINA = RGBColor(0xFF, 0x7A, 0x45)    # Orange
COLOR_MEDIA = RGBColor(0xAB, 0x47, 0xBC)   # Purple
COLOR_OURS = RGBColor(0x66, 0xBB, 0x6A)    # Green

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_card_with_title(slide, left, top, width, height, title, title_color=ACCENT_BLUE):
    card = add_rounded_rect(slide, left, top, width, height)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
                 title, font_size=16, color=title_color, bold=True)
    return card


def add_table(slide, left, top, width, height, rows, cols, data, header_color=ACCENT_BLUE):
    """Add a styled table. data is list of lists [row][col]."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(data[i][j])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = 'Arial'
                paragraph.alignment = PP_ALIGN.CENTER
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = TEXT_WHITE
                else:
                    paragraph.font.color.rgb = TEXT_WHITE
                    paragraph.font.bold = False

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4A)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x38) if i % 2 == 1 else RGBColor(0x28, 0x28, 0x42)

    return table_shape


def add_multiline_text(slide, left, top, width, height, lines, font_size=14, color=TEXT_WHITE,
                       line_colors=None, line_bolds=None, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        prefix = "  " if bullet else ""
        p.text = prefix + line
        p.font.size = Pt(font_size)
        p.font.name = 'Arial'
        p.font.color.rgb = (line_colors[i] if line_colors and i < len(line_colors) else color)
        p.font.bold = (line_bolds[i] if line_bolds and i < len(line_bolds) else False)
        p.space_after = Pt(4)

    return txBox


# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
             "Baseline Algorithm Correctness Audit", font_size=40, color=TEXT_WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
             "SGX TEE Edge Cluster DNN Inference Scheduling", font_size=22, color=ACCENT_BLUE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             "v7: OCC / DINA / MEDIA correctness audit against original papers",
             font_size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "2026-03-09", font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "Agenda", font_size=32, color=TEXT_WHITE, bold=True)

items = [
    ("1.", "Baseline Fix Summary", "OCC / DINA / MEDIA correctness deviations & corrections"),
    ("2.", "Exp1: Fixed Comparison", "4x Xeon, 100Mbps - all 12 models"),
    ("3.", "Exp2: Network Ablation", "Bandwidth sweep 0.5-500 Mbps"),
    ("4.", "Exp3: Server Ablation", "Heterogeneous 1-8 servers"),
    ("5.", "Key Findings", "OCC=DINA=MEDIA phenomenon & Ours advantage analysis"),
    ("6.", "Next Steps", "Remaining work items"),
]

for i, (num, title, desc) in enumerate(items):
    y = Inches(1.3 + i * 0.95)
    add_rounded_rect(slide, Inches(1), y, Inches(11), Inches(0.8))
    add_text_box(slide, Inches(1.3), y + Inches(0.08), Inches(0.5), Inches(0.35),
                 num, font_size=20, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.9), y + Inches(0.08), Inches(4), Inches(0.35),
                 title, font_size=18, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(1.9), y + Inches(0.42), Inches(9), Inches(0.3),
                 desc, font_size=13, color=TEXT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 3: Baseline Fix Summary
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "1. Baseline Algorithm Corrections (commit ead9b1b)", font_size=28, color=TEXT_WHITE, bold=True)

# OCC card
add_card_with_title(slide, Inches(0.5), Inches(1.2), Inches(4), Inches(2.7),
                    "OCC (Occlumency, MobiCom'19)", ACCENT_BLUE)
add_multiline_text(slide, Inches(0.7), Inches(1.7), Inches(3.5), Inches(2.0), [
    "Severity: Medium",
    "",
    "Bug: Weights included in EPC budget",
    "Paper: Weights stored outside EPC",
    "  (unprotected DRAM, DDR memcpy)",
    "",
    "Fix: activation-only EPC budget",
    "  EPC = activations + ring buffer (20MB)",
    "  Budget = 93 - 20 = 73 MB",
], font_size=12, color=TEXT_WHITE,
    line_colors=[ACCENT_ORANGE, TEXT_WHITE, ACCENT_RED, ACCENT_GREEN, ACCENT_GREEN,
                 TEXT_WHITE, ACCENT_GREEN, TEXT_GRAY, TEXT_GRAY])

# DINA card
add_card_with_title(slide, Inches(4.7), Inches(1.2), Inches(4), Inches(2.7),
                    "DINA (IEEE TPDS'24)", ACCENT_ORANGE)
add_multiline_text(slide, Inches(4.9), Inches(1.7), Inches(3.5), Inches(2.0), [
    "Severity: HIGH",
    "",
    "Bug: Forced server rotation",
    "  (round-robin partition assignment)",
    "Paper: Matching game approach",
    "  DINA-P: adaptive partitioning",
    "  DINA-O: greedy + pairwise swap",
    "",
    "Fix: Full rewrite (319 lines)",
    "  Proportional workload distribution",
], font_size=12, color=TEXT_WHITE,
    line_colors=[ACCENT_RED, TEXT_WHITE, ACCENT_RED, ACCENT_RED, ACCENT_GREEN, ACCENT_GREEN,
                 ACCENT_GREEN, TEXT_WHITE, ACCENT_GREEN, TEXT_GRAY])

# MEDIA card
add_card_with_title(slide, Inches(8.9), Inches(1.2), Inches(4), Inches(2.7),
                    "MEDIA (ICDCS)", ACCENT_PURPLE)
add_multiline_text(slide, Inches(9.1), Inches(1.7), Inches(3.5), Inches(2.0), [
    "Severity: Low-Medium",
    "",
    "Bug: Constraint 1 too strict",
    "  Only checked in_degree(v)==1",
    "Paper: in_deg==1 OR out_deg==1",
    "",
    "Fix: Restored paper's Constraint 1",
    "  + Priority computation (Eq. 11)",
    "  rewritten to match paper formula",
], font_size=12, color=TEXT_WHITE,
    line_colors=[ACCENT_ORANGE, TEXT_WHITE, ACCENT_RED, ACCENT_RED, ACCENT_GREEN,
                 TEXT_WHITE, ACCENT_GREEN, TEXT_GRAY, TEXT_GRAY])

# Impact summary
add_rounded_rect(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(3.0), RGBColor(0x1E, 0x1E, 0x35))
add_text_box(slide, Inches(0.7), Inches(4.3), Inches(11.5), Inches(0.4),
             "Impact on Results", font_size=18, color=ACCENT_BLUE, bold=True)
add_multiline_text(slide, Inches(0.7), Inches(4.8), Inches(11.5), Inches(2.2), [
    "DINA (biggest change): Old code forced round-robin assignment -> catastrophic latency (e.g., BERT-large: 17,532ms)",
    "  New code: greedy scheduler assigns all partitions to best server -> matches OCC/MEDIA on serial DAGs",
    "",
    "Result: OCC = DINA = MEDIA on ALL 12 models (homogeneous 4x Xeon servers)",
    "  This is CORRECT - proves serial DAGs cannot benefit from multi-server distribution",
    "  Only Ours (HPA tensor parallelism) can exploit intra-layer parallelism",
], font_size=13, color=TEXT_WHITE,
    line_colors=[TEXT_WHITE, TEXT_GRAY, TEXT_WHITE, ACCENT_GREEN, TEXT_GRAY, ACCENT_BLUE],
    line_bolds=[True, False, False, True, False, True])


# ════════════════════════════════════════════════════════════════
# SLIDE 4: Exp1 Results
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "2. Exp1: Fixed Configuration (4x Xeon_IceLake, 100 Mbps)", font_size=28, color=TEXT_WHITE, bold=True)

exp1_data = [
    ["Model", "OCC", "DINA", "MEDIA", "Ours", "Improvement"],
    ["InceptionV3", "1506", "1506", "1506", "930", "38.3%"],
    ["ViT-base", "1218", "1218", "1218", "1046", "14.1%"],
    ["ViT-small", "410", "410", "410", "334", "18.4%"],
    ["BERT-base", "757", "757", "757", "663", "12.4%"],
    ["ALBERT-base", "756", "756", "756", "660", "12.8%"],
    ["DistilBERT", "377", "377", "377", "337", "10.5%"],
    ["TinyBERT-4l", "81", "81", "81", "75", "7.2%"],
    ["TinyBERT-6l", "377", "377", "377", "337", "10.5%"],
    ["ViT-tiny", "180", "180", "180", "160", "11.3%"],
    ["BERT-large", "2307", "2307", "2307", "2287", "0.9%"],
    ["ALBERT-large", "2382", "2382", "2382", "2288", "4.0%"],
    ["ViT-large", "3563", "3563", "3563", "3526", "1.0%"],
]

add_table(slide, Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.6),
          len(exp1_data), 6, exp1_data)

# Annotations
add_card_with_title(slide, Inches(9.0), Inches(1.1), Inches(4.0), Inches(5.6),
                    "Key Observations", ACCENT_BLUE)
add_multiline_text(slide, Inches(9.2), Inches(1.6), Inches(3.6), Inches(5.0), [
    "OCC = DINA = MEDIA",
    "All 12 models, exact same latency",
    "",
    "Why? Serial DAG chains cannot",
    "benefit from multi-server distribution.",
    "All 3 baselines effectively run on",
    "a single best server.",
    "",
    "Ours advantage by model type:",
    "",
    "  Parallel branches (InceptionV3):",
    "  38.3% - largest gain",
    "",
    "  Small models (ViT-s/t, BERT-b):",
    "  12-18% - TP Amdahl speedup",
    "",
    "  Large models (BERT-l, ViT-l):",
    "  ~1% - paging penalty offsets TP",
], font_size=12, color=TEXT_WHITE,
    line_colors=[ACCENT_GREEN, TEXT_GRAY, TEXT_WHITE, TEXT_WHITE, TEXT_WHITE, TEXT_WHITE, TEXT_WHITE,
                 TEXT_WHITE,
                 ACCENT_BLUE, TEXT_WHITE,
                 ACCENT_GREEN, ACCENT_GREEN, TEXT_WHITE,
                 ACCENT_BLUE, TEXT_GRAY, TEXT_WHITE,
                 ACCENT_ORANGE, TEXT_GRAY],
    line_bolds=[True, False, False, False, False, False, False, False, True])


# ════════════════════════════════════════════════════════════════
# SLIDE 5: Exp2 - InceptionV3 Network Ablation
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "3. Exp2: Network Bandwidth Ablation (4 servers)", font_size=28, color=TEXT_WHITE, bold=True)

# InceptionV3 sub-table
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(6), Inches(0.4),
             "InceptionV3 (parallel branch model - key result)", font_size=16, color=ACCENT_GREEN, bold=True)

iv3_data = [
    ["BW (Mbps)", "OCC", "DINA", "MEDIA", "Ours", "Ours/OCC"],
    ["0.5", "1506", "1506", "1506", "1506", "1.000x"],
    ["5.0", "1506", "1506", "1506", "1499", "0.996x"],
    ["10", "1506", "1506", "1506", "1308", "0.869x"],
    ["50", "1506", "1506", "1506", "1000", "0.664x"],
    ["100", "1506", "1506", "1506", "930", "0.617x"],
    ["200", "1506", "1506", "1506", "867", "0.576x"],
    ["500", "1506", "1506", "1506", "556", "0.370x"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(7.5), Inches(3.4),
          len(iv3_data), 6, iv3_data)

# BERT-large sub-table
add_text_box(slide, Inches(0.5), Inches(5.1), Inches(6), Inches(0.4),
             "BERT-large (large linear model)", font_size=16, color=ACCENT_BLUE, bold=True)

bl_data = [
    ["BW (Mbps)", "OCC", "DINA", "MEDIA", "Ours", "Ours/OCC"],
    ["0.5-200", "2307", "2307", "2307", "2287", "0.991x"],
    ["500", "2307", "2307", "2408", "1817", "0.787x"],
]
add_table(slide, Inches(0.3), Inches(5.5), Inches(7.5), Inches(1.3),
          len(bl_data), 6, bl_data)

# Right side annotations
add_card_with_title(slide, Inches(8.2), Inches(1.0), Inches(4.8), Inches(6.0),
                    "Bandwidth Sensitivity Analysis", ACCENT_BLUE)
add_multiline_text(slide, Inches(8.4), Inches(1.5), Inches(4.4), Inches(5.3), [
    "InceptionV3: S-curve pattern",
    "",
    "  <5 Mbps: safeguard (=OCC)",
    "    Comm cost > parallel benefit",
    "",
    "  10 Mbps: 13% speedup (HEFT only)",
    "  100 Mbps: 38% speedup (HEFT+TP)",
    "  500 Mbps: 63% speedup (max TP)",
    "",
    "  TP operators increase with BW:",
    "  10Mbps: 1 op, 100Mbps: 17 ops",
    "  500Mbps: 58 ops (all at k=8)",
    "",
    "BERT-large: safeguard dominant",
    "",
    "  Residual connections prevent HEFT",
    "  500 Mbps breakthrough: 21% speedup",
    "  (24 ops get TP, breaks dependencies)",
    "",
    "Baselines bandwidth-insensitive",
    "  All 3 always use single server",
], font_size=12, color=TEXT_WHITE,
    line_colors=[ACCENT_GREEN, TEXT_WHITE,
                 TEXT_GRAY, TEXT_DIM, TEXT_WHITE,
                 TEXT_WHITE, TEXT_WHITE, ACCENT_GREEN, TEXT_WHITE,
                 ACCENT_BLUE, TEXT_GRAY, TEXT_GRAY, TEXT_WHITE,
                 ACCENT_ORANGE, TEXT_WHITE,
                 TEXT_GRAY, TEXT_WHITE, TEXT_DIM, TEXT_WHITE,
                 ACCENT_RED, TEXT_GRAY],
    line_bolds=[True, False, False, False, False, False, False, True, False,
                False, False, False, False, True, False, False, False, False, False, True])


# ════════════════════════════════════════════════════════════════
# SLIDE 6: Exp3 - Server Ablation
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "4. Exp3: Heterogeneous Server Ablation (100 Mbps)", font_size=28, color=TEXT_WHITE, bold=True)

# Server order description
add_text_box(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(0.4),
             "Server addition order: 2x Celeron(0.11x) + 4x i5-6500(0.93x) + i3-10100(1.03x) + i5-11600(1.97x)",
             font_size=13, color=TEXT_GRAY)

# InceptionV3 table
add_text_box(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(0.3),
             "InceptionV3", font_size=16, color=ACCENT_GREEN, bold=True)

iv3_server_data = [
    ["Servers", "Config", "OCC", "DINA", "MEDIA", "Ours", "Ours Imprv."],
    ["1", "1x Celeron", "13689", "13689", "13689", "13689", "0%"],
    ["2", "2x Celeron", "13689", "13689", "13689", "9887", "27.8%"],
    ["3", "+i5-6500", "1619", "1619", "1619", "1472", "9.1%"],
    ["4", "+i5-6500", "1619", "1619", "1619", "1168", "27.8%"],
    ["5-6", "+i5-6500s", "1619", "1619", "1619", "986-1049", "35-39%"],
    ["7", "+i3-10100", "1462", "1462", "1462", "872", "40.4%"],
    ["8", "+i5-11600", "764", "764", "764", "511", "33.2%"],
]
add_table(slide, Inches(0.3), Inches(1.8), Inches(8.5), Inches(3.3),
          len(iv3_server_data), 7, iv3_server_data)

# BERT-base table
add_text_box(slide, Inches(0.5), Inches(5.3), Inches(6), Inches(0.3),
             "BERT-base (linear model reference)", font_size=16, color=ACCENT_BLUE, bold=True)

bert_server_data = [
    ["Servers", "OCC", "DINA", "MEDIA", "Ours", "Ours/OCC"],
    ["1", "6883", "6883", "6883", "6883", "1.000x"],
    ["2", "6883", "6883", "6554", "6261", "0.909x"],
    ["3-6", "814", "814", "814", "713", "0.876x"],
    ["8", "384", "384", "384", "337", "0.877x"],
]
add_table(slide, Inches(0.3), Inches(5.7), Inches(7.5), Inches(1.5),
          len(bert_server_data), 6, bert_server_data)

# Right annotations
add_card_with_title(slide, Inches(9.0), Inches(1.4), Inches(4.0), Inches(5.8),
                    "Server Scaling Analysis", ACCENT_BLUE)
add_multiline_text(slide, Inches(9.2), Inches(1.9), Inches(3.6), Inches(5.0), [
    "Key transition points:",
    "",
    "  n=1->2 (Celeron 0.11x):",
    "  Baselines: no change (too slow)",
    "  Ours: 28% speedup via HPA",
    "  -> HPA works even on slow nodes!",
    "",
    "  n=2->3 (add i5-6500 0.93x):",
    "  All methods: 13689 -> 1619 ms",
    "  (8.5x drop, first usable server)",
    "",
    "  n=7 (best Ours/OCC = 40.4%):",
    "  Peak improvement point",
    "",
    "Baselines can only benefit from",
    "faster single servers (OCC jumps),",
    "never from more servers.",
    "",
    "Ours: sub-linear scaling with",
    "diminishing returns after n=5.",
], font_size=12, color=TEXT_WHITE,
    line_colors=[ACCENT_BLUE, TEXT_WHITE,
                 TEXT_WHITE, TEXT_GRAY, ACCENT_GREEN, ACCENT_GREEN, TEXT_WHITE,
                 TEXT_WHITE, TEXT_GRAY, TEXT_DIM, TEXT_WHITE,
                 ACCENT_GREEN, TEXT_GRAY, TEXT_WHITE,
                 TEXT_WHITE, TEXT_WHITE, TEXT_WHITE, TEXT_WHITE,
                 ACCENT_BLUE, TEXT_GRAY],
    line_bolds=[True, False, False, False, True, False, False, False, False, False, False, True])


# ════════════════════════════════════════════════════════════════
# SLIDE 7: Key Findings Summary
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "5. Key Findings & Thesis Argument", font_size=28, color=TEXT_WHITE, bold=True)

# Finding 1
add_card_with_title(slide, Inches(0.5), Inches(1.2), Inches(6.0), Inches(2.5),
                    "Finding 1: OCC = DINA = MEDIA", ACCENT_GREEN)
add_multiline_text(slide, Inches(0.7), Inches(1.7), Inches(5.5), Inches(2.0), [
    "All 3 baselines produce identical latency",
    "on homogeneous servers (all 12 models).",
    "",
    "Root cause: Serial DAG chains prevent",
    "multi-server parallelism. Partition-level",
    "distribution adds communication overhead",
    "without parallel execution benefit.",
    "",
    "This validates the thesis core argument.",
], font_size=13, color=TEXT_WHITE,
    line_colors=[TEXT_WHITE, TEXT_WHITE, TEXT_WHITE, TEXT_WHITE, TEXT_WHITE, TEXT_WHITE, TEXT_WHITE,
                 TEXT_WHITE, ACCENT_GREEN],
    line_bolds=[True, False, False, False, False, False, False, False, True])

# Finding 2
add_card_with_title(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(2.5),
                    "Finding 2: Ours advantage depends on model topology", ACCENT_BLUE)
add_multiline_text(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(2.0), [
    "Parallel branches (InceptionV3): 38-63%",
    "  -> HPA tensor parallelism + HEFT scheduling",
    "",
    "Small linear models: 12-18%",
    "  -> Pure TP Amdahl speedup (k=8)",
    "",
    "Large linear models: 1% @ 100Mbps",
    "  -> Paging penalty dominates",
    "  -> But 500Mbps: 21-23% breakthrough!",
], font_size=13, color=TEXT_WHITE,
    line_colors=[ACCENT_GREEN, TEXT_GRAY, TEXT_WHITE, ACCENT_BLUE, TEXT_GRAY, TEXT_WHITE,
                 ACCENT_ORANGE, TEXT_GRAY, ACCENT_GREEN],
    line_bolds=[True, False, False, True, False, False, True])

# Finding 3
add_card_with_title(slide, Inches(0.5), Inches(4.0), Inches(6.0), Inches(3.0),
                    "Finding 3: Bandwidth-adaptive behavior", ACCENT_PURPLE)
add_multiline_text(slide, Inches(0.7), Inches(4.5), Inches(5.5), Inches(2.5), [
    "Ours exhibits S-curve speedup pattern:",
    "",
    "  Low BW (<5Mbps): safeguard = OCC",
    "    No degradation (unlike old DINA: 100x+)",
    "",
    "  Medium BW (10-100Mbps): HEFT + progressive TP",
    "    Operators selected by DP based on BW",
    "",
    "  High BW (500Mbps): maximum TP utilization",
    "    InceptionV3: 63%, BERT-large: 21%",
    "",
    "Adaptive degradation = safety guarantee",
], font_size=13, color=TEXT_WHITE,
    line_colors=[TEXT_WHITE, TEXT_WHITE, TEXT_GRAY, ACCENT_GREEN, TEXT_WHITE,
                 TEXT_WHITE, TEXT_GRAY, TEXT_WHITE, TEXT_WHITE, ACCENT_GREEN,
                 TEXT_WHITE, ACCENT_BLUE],
    line_bolds=[True, False, False, False, False, False, False, False, False, True, False, True])

# Finding 4
add_card_with_title(slide, Inches(6.8), Inches(4.0), Inches(6.0), Inches(3.0),
                    "Finding 4: Heterogeneous server utilization", ACCENT_ORANGE)
add_multiline_text(slide, Inches(7.0), Inches(4.5), Inches(5.5), Inches(2.5), [
    "Baselines: only benefit from faster servers",
    "  Cannot exploit multiple servers for 1 inference",
    "  OCC jumps at n=3 (first fast server added)",
    "",
    "Ours: progressive improvement with more servers",
    "  Even slow Celeron (0.11x) contributes via HPA",
    "  Sub-linear scaling (diminishing returns >5)",
    "",
    "InceptionV3 n=7: Ours = 872ms (40% < baseline)",
    "  vs baselines stuck at 1462ms",
    "",
    "Demonstrates true multi-server utilization",
], font_size=13, color=TEXT_WHITE,
    line_colors=[TEXT_WHITE, TEXT_GRAY, TEXT_GRAY, TEXT_WHITE,
                 ACCENT_GREEN, TEXT_GRAY, TEXT_GRAY, TEXT_WHITE,
                 ACCENT_GREEN, TEXT_GRAY, TEXT_WHITE, ACCENT_BLUE],
    line_bolds=[True, False, False, False, True, False, False, False, True, False, False, True])


# ════════════════════════════════════════════════════════════════
# SLIDE 8: Version History (brief)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "Version History: Bug Fix Journey (v1 -> v7)", font_size=28, color=TEXT_WHITE, bold=True)

versions = [
    ["Version", "Date", "Fix Description", "Impact"],
    ["v1", "03-03", "Initial implementation", "Baseline (comm volume bug)"],
    ["v2", "03-03", "AllReduce + safeguard fix", "Non-deterministic set() bug"],
    ["v3", "03-04", "Non-determinism fix (sorted sets)", "Reproducible results"],
    ["v4", "03-04", "OCC: EPC paging -> DDR memcpy", "OCC baseline 64-81% lower"],
    ["v5", "03-05", "MEDIA: remove fake paging_cost", "MEDIA = OCC (correct)"],
    ["v6", "03-05", "MEDIA: fix join node merging", "MEDIA < OCC on InceptionV3"],
    ["v7", "03-09", "OCC/DINA/MEDIA full audit", "OCC=DINA=MEDIA (correct!)"],
]
add_table(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.5),
          len(versions), 4, versions)

add_text_box(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
             "7 versions, 6+ bugs fixed. Each fix brought results closer to ground truth.",
             font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(1.5), [
    "Key lesson: Every baseline implementation had deviations from original papers.",
    "Rigorous audit against source PDFs was essential for credible experimental comparison.",
    "The final v7 result (OCC=DINA=MEDIA) is actually the strongest thesis argument -",
    "proving that ALL existing partition-based methods fail to exploit multi-server parallelism.",
], font_size=14, color=TEXT_WHITE,
    line_bolds=[True, False, False, True])


# ════════════════════════════════════════════════════════════════
# SLIDE 9: Next Steps
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "6. Next Steps", font_size=28, color=TEXT_WHITE, bold=True)

# TODO items
todos = [
    ("HIGH", "Re-run Exp2 for remaining models", "Exp2 was partially stale (segfault during generation). BERT-large/ALBERT-large/ViT-large etc. need fresh data with fixed algorithms.", ACCENT_RED),
    ("HIGH", "Investigate ALBERT-large segfault", "Exp2 crashed at ALBERT-large (1371 nodes, exit code 139). Need per-model subprocess isolation.", ACCENT_RED),
    ("MED", "Re-generate all Exp2 charts", "After fresh data, regenerate network ablation line charts and combined grid.", ACCENT_ORANGE),
    ("MED", "Paper writing: Evaluation section", "Organize v7 results into paper format. Key story: OCC=DINA=MEDIA proves baseline limitation.", ACCENT_ORANGE),
    ("LOW", "Additional ablation experiments", "P_sync sensitivity, HPA k-distribution breakdown charts for paper.", ACCENT_BLUE),
    ("LOW", "Patent application", "Draft CN patent for HPA tensor parallelism + HEFT scheduling method.", ACCENT_BLUE),
]

for i, (priority, title, desc, color) in enumerate(todos):
    y = Inches(1.2 + i * 1.0)
    add_rounded_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.9))

    # Priority badge
    badge = add_rounded_rect(slide, Inches(0.7), y + Inches(0.15), Inches(0.7), Inches(0.3), color)
    add_text_box(slide, Inches(0.7), y + Inches(0.15), Inches(0.7), Inches(0.3),
                 priority, font_size=10, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.6), y + Inches(0.1), Inches(10), Inches(0.35),
                 title, font_size=16, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(1.6), y + Inches(0.48), Inches(10.5), Inches(0.35),
                 desc, font_size=12, color=TEXT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 10: Thank you
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "Thank You", font_size=44, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             "Questions & Discussion", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(1.5),
             "Key takeaway: After rigorous baseline audit, OCC=DINA=MEDIA on all serial DAGs.\n"
             "Only HPA (Ours) achieves real multi-server parallelism through tensor parallelism.",
             font_size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ──
output_path = os.path.join(os.path.dirname(__file__), "group_meeting_v7_report.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
